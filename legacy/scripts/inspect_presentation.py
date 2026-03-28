from pathlib import Path

from pptx import Presentation


def main() -> None:
    path = Path("Semantic_Map_BoE.pptx")
    presentation = Presentation(path)
    print(f"slides={len(presentation.slides)}")
    print(f"layouts={len(presentation.slide_layouts)}")

    for index, layout in enumerate(presentation.slide_layouts):
        print(f"LAYOUT {index}: {layout.name}")
        for shape in layout.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                print(f"  {text.replace(chr(10), ' | ')}")

    print("--- slides ---")
    for index, slide in enumerate(presentation.slides, start=1):
        print(f"SLIDE {index}")
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                print(f"  {text.replace(chr(10), ' | ')}")


if __name__ == "__main__":
    main()

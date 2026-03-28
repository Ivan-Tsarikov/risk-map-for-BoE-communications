from __future__ import annotations

import csv
import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE_PATH = ROOT / "presentation" / "Semantic_Map_BoE_template.pptx"
OUTPUT_PATH = ROOT / "presentation" / "Semantic_Map_BoE_public.pptx"
FALLBACK_OUTPUT_PATH = ROOT / "presentation" / "Semantic_Map_BoE_public_updated.pptx"

REPORT_PATH = ROOT / "data" / "output" / "public_cleaning_summary.json"
HOT_PATH = ROOT / "data" / "output" / "hot_regions_topics.csv"
FIGURES_DIR = ROOT / "reports" / "figures"

SLIDE_W = Cm(33.867)
SLIDE_H = Cm(19.05)
PAGE_MARGIN_X = Cm(1.2)

TEXT_DARK = RGBColor(34, 34, 34)
TEXT_MUTED = RGBColor(95, 95, 95)
ACCENT_BLUE = RGBColor(37, 86, 167)
ACCENT_RED = RGBColor(166, 38, 46)
BOX_FILL = RGBColor(247, 248, 250)
BOX_LINE = RGBColor(210, 214, 220)


def load_presentation() -> Presentation:
    if TEMPLATE_PATH.exists():
        presentation = Presentation(TEMPLATE_PATH)
    else:
        presentation = Presentation()
        presentation.slide_width = SLIDE_W
        presentation.slide_height = SLIDE_H
    presentation.slide_width = SLIDE_W
    presentation.slide_height = SLIDE_H
    return presentation


def remove_all_slides(presentation: Presentation) -> None:
    slide_ids = list(presentation.slides._sldIdLst)
    for slide_id in slide_ids:
        relationship_id = slide_id.rId
        presentation.part.drop_rel(relationship_id)
        presentation.slides._sldIdLst.remove(slide_id)


def blank_layout(presentation: Presentation):
    return presentation.slide_layouts[6] if len(presentation.slide_layouts) > 6 else presentation.slide_layouts[-1]


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Cm(1.2), Cm(0.8), Cm(31.2), Cm(2.4))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Cm(1.2), Cm(2.7), Cm(31.2), Cm(1.4))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.LEFT


def add_paragraph_box(
    slide,
    text_blocks: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 14,
    title: str | None = None,
    title_color: RGBColor | None = None,
    fill: bool = False,
) -> None:
    if fill:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = BOX_FILL
        shape.line.color.rgb = BOX_LINE
        text_frame = shape.text_frame
    else:
        box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = box.text_frame

    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    if title:
        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(font_size + 1)
        p.font.bold = True
        p.font.color.rgb = title_color or ACCENT_RED
        p.space_after = Pt(8)
        start_index = 0
    else:
        start_index = -1

    for idx, block in enumerate(text_blocks):
        p = text_frame.paragraphs[0] if idx == 0 and start_index == -1 else text_frame.add_paragraph()
        p.text = block
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(7)
        p.alignment = PP_ALIGN.LEFT


def add_metric_card(slide, title: str, value: str, left: float, top: float, width: float, height: float) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BOX_FILL
    shape.line.color.rgb = BOX_LINE
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = frame.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(11)
    p1.font.color.rgb = TEXT_MUTED
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(5)

    p2 = frame.add_paragraph()
    p2.text = value
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_BLUE
    p2.alignment = PP_ALIGN.CENTER


def add_table(slide, rows: list[list[str]], left: float, top: float, width: float, height: float) -> None:
    table = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height).table
    col_width = int(width / len(rows[0]))
    for col_idx in range(len(rows[0])):
        table.columns[col_idx].width = col_width

    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = TEXT_DARK
                    if row_idx == 0:
                        run.font.bold = True


def add_picture_fit(slide, image_path: Path, left: float, top: float, width: float, height: float) -> None:
    if not image_path.exists():
        add_paragraph_box(
            slide,
            [f"Файл не найден: {image_path.name}"],
            left,
            top,
            width,
            height,
            font_size=12,
            fill=True,
        )
        return

    with Image.open(image_path) as image:
        image_ratio = image.width / image.height
    box_ratio = width / height

    if image_ratio > box_ratio:
        pic_width = width
        pic_height = width / image_ratio
        pic_left = left
        pic_top = top + (height - pic_height) / 2
    else:
        pic_height = height
        pic_width = height * image_ratio
        pic_top = top
        pic_left = left + (width - pic_width) / 2

    slide.shapes.add_picture(str(image_path), pic_left, pic_top, width=pic_width, height=pic_height)


def format_int(value: int) -> str:
    return f"{int(value):,}".replace(",", " ")


def read_cleaning_summary() -> dict:
    with REPORT_PATH.open("r", encoding="utf-8") as handle:
        return json.load(handle)


def read_hot_pairs(limit: int = 6) -> list[dict[str, str]]:
    if not HOT_PATH.exists():
        return []
    with HOT_PATH.open("r", encoding="utf-8-sig", newline="") as handle:
        rows = list(csv.DictReader(handle))
    return rows[:limit]


def build_slide_1_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title(
        slide,
        "Semantic Map of Regional Media Agendas for Central-Bank Communication",
        "MVP системы раннего выявления региональных коммуникационных рисков на локальных новостях",
    )
    add_paragraph_box(
        slide,
        [
            "Проект строит интерпретируемую semantic map региональных тем, значимых для коммуникации центрального банка.",
            "Основной результат — недельный региональный индекс, который показывает, где конкретная тема усиливается относительно собственной локальной истории.",
        ],
        Cm(1.2),
        Cm(5.0),
        Cm(16.5),
        Cm(5.2),
        font_size=16,
    )
    add_paragraph_box(
        slide,
        [
            "Pipeline: news corpus -> candidate gate -> story embeddings -> econ gate -> BoE topics -> regional risk index -> maps and trends",
            "Текущая версия результатов относится к режиму sampled semantic pilot.",
        ],
        Cm(1.2),
        Cm(11.0),
        Cm(16.5),
        Cm(3.8),
        font_size=13,
        fill=True,
        title="Краткая схема",
        title_color=ACCENT_BLUE,
    )
    add_picture_fit(slide, FIGURES_DIR / "overall_top_risk_map.png", Cm(18.8), Cm(4.8), Cm(13.4), Cm(12.2))


def build_slide_2_problem(prs: Presentation) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title(slide, "Постановка задачи", "Задача проекта — перейти от локальных новостей к интерпретируемой региональной карте тем и рисков.")
    add_paragraph_box(
        slide,
        [
            "Центральному банку важно видеть не только общенациональную медиаповестку, но и региональные темы, которые усиливаются раньше или заметнее в локальном информационном поле.",
            "Локальные онлайн-медиа позволяют наблюдать эту динамику на уровне административных районов, если тексты очищены, снабжены региональной привязкой и приведены к единому временно́му шагу.",
            "Цель MVP состоит в том, чтобы для каждой недели определить, какая тема Bank of England усиливается в конкретном регионе и насколько этот сигнал отклоняется от собственной локальной нормы.",
        ],
        Cm(1.2),
        Cm(4.5),
        Cm(15.4),
        Cm(10.4),
        font_size=15,
    )
    add_paragraph_box(
        slide,
        [
            "Этап 1. Очистка и нормализация корпуса.",
            "Этап 2. Отбор candidate-статей.",
            "Этап 3. Story-level embeddings.",
            "Этап 4. Semantic econ gate и присвоение тем BoE.",
            "Этап 5. Построение регионального индекса, карт и трендов.",
        ],
        Cm(18.2),
        Cm(5.2),
        Cm(13.6),
        Cm(9.2),
        font_size=14,
        fill=True,
        title="Логика решения",
        title_color=ACCENT_BLUE,
    )


def build_slide_3_dataset_source(prs: Presentation) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title(slide, "Источник данных и пригодность корпуса", "Источник данных — articles.csv из открытого корпуса UKTwitNewsCor.")
    add_paragraph_box(
        slide,
        [
            "Корпус включает более 2,5 млн статей локальных онлайн-медиа Великобритании за январь 2020 — декабрь 2022, собранных по 360 локальным доменам и дополненных региональными и social-media метаданными.",
            "UKTwitNewsCor покрывает 94% Local Authority Districts и позволяет анализировать локальные медиаповестки по LAD / main_LAD.",
            "Для проекта это важно, потому что корпус сочетает тексты новостей, временные метки, ссылки, метрики вовлечённости и региональную привязку.",
            "Методологически проект измеряет не точную географию событий, а структуру локального медиаполя по административным районам.",
        ],
        Cm(1.2),
        Cm(4.6),
        Cm(19.1),
        Cm(10.6),
        font_size=15,
    )
    add_paragraph_box(
        slide,
        [
            "Ключевые поля корпуса: article_text, article_date, tweet_date, resolved_url, LAD, main_LAD, engagement fields.",
            "Региональная единица анализа: LAD / main_LAD.",
            "Публичная версия репозитория не включает сам корпус и работает только с code-first и licence-safe результатами.",
        ],
        Cm(21.0),
        Cm(5.1),
        Cm(10.7),
        Cm(8.6),
        font_size=13,
        fill=True,
        title="Почему корпус подходит",
        title_color=ACCENT_BLUE,
    )


def build_slide_4_dataset_coverage(prs: Presentation, summary: dict) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title(slide, "Структура корпуса и покрытие", "Ниже показаны основные численные характеристики очищенного корпуса.")
    add_metric_card(slide, "Сырые строки", format_int(summary["rows_seen"]), Cm(1.2), Cm(4.4), Cm(6.1), Cm(3.0))
    add_metric_card(slide, "Строки после очистки", format_int(summary["rows_kept"]), Cm(7.7), Cm(4.4), Cm(6.1), Cm(3.0))
    add_metric_card(slide, "Исключено как слишком короткие", format_int(summary["rows_dropped_short_text"]), Cm(14.2), Cm(4.4), Cm(6.1), Cm(3.0))
    add_metric_card(slide, "Исключено по дате", format_int(summary["rows_dropped_date_outlier"]), Cm(20.7), Cm(4.4), Cm(6.1), Cm(3.0))
    add_metric_card(slide, "Дубликатные story groups", format_int(summary["duplicate_story_groups"]), Cm(27.2), Cm(4.4), Cm(5.4), Cm(3.0))

    add_paragraph_box(
        slide,
        [
            summary.get("source_period_note", "Источник корпуса по статье: январь 2020 — декабрь 2022."),
            summary.get("semantic_run_note", "Текущий semantic run относится к sampled semantic pilot."),
            "Полные локальные parquet-артефакты и embeddings не публикуются в GitHub.",
        ],
        Cm(1.2),
        Cm(8.2),
        Cm(31.0),
        Cm(3.2),
        font_size=14,
        fill=True,
        title="Интерпретация",
        title_color=ACCENT_BLUE,
    )
    add_paragraph_box(
        slide,
        [
            "После очистки для каждой строки нормализуются дата, регион, очищенный текст, неделя и story_key.",
            "Story_key строится иерархически: duplicate_group -> article_id -> unique_article_id. Это позволяет считать семантику один раз на историю, а не на каждую перепубликацию.",
        ],
        Cm(1.2),
        Cm(12.0),
        Cm(31.0),
        Cm(4.2),
        font_size=14,
    )


def build_slide_5_cleaning(prs: Presentation) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title(slide, "Как очищались тексты и готовился корпус", "Очистка стандартизирует текст, дату и регион ещё до semantic-этапа.")
    add_paragraph_box(
        slide,
        [
            "Шаг 1. Дата статьи берётся из article_date, а при её отсутствии или низком доверии используется tweet_date.",
            "Шаг 2. Регион нормализуется по main_LAD, а при его отсутствии по LAD.",
            "Шаг 3. Текст очищается от лишних пробелов, слишком короткие тексты удаляются.",
            "Шаг 4. Для всех строк создаётся week и укороченное поле text_for_embedding, которое содержит первые 1200 символов cleaned text.",
            "Шаг 5. Для каждой истории формируется story_key и вес dup_weight, чтобы перепубликации не завышали частоты.",
        ],
        Cm(1.2),
        Cm(4.7),
        Cm(15.3),
        Cm(10.8),
        font_size=14,
    )
    add_paragraph_box(
        slide,
        [
            "dup_weight = 1 / count(story_key, region, week)",
            "Если одна и та же история опубликована в одном регионе и в одну неделю несколько раз, её вклад делится между дубликатами.",
            "Эта нормализация работает ещё до построения регионального индекса и снижает влияние механических перепубликаций.",
        ],
        Cm(18.0),
        Cm(5.2),
        Cm(13.8),
        Cm(6.8),
        font_size=14,
        fill=True,
        title="Зачем нужен dup_weight",
        title_color=ACCENT_BLUE,
    )


def build_slide_6_econ_gate(prs: Presentation) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title(slide, "Как отбирались экономические статьи", "Economic relevance задаётся отдельным шлюзом и не сводится к 7 темам BoE.")
    add_paragraph_box(
        slide,
        [
            "Сначала строится candidate set. В него попадают статьи с strong keyword hit, weak keyword hit или url_section hit.",
            "После этого семантика считается на уровне story_key. В embedding подаётся не полный текст, а поле text_for_embedding — первые 1200 символов очищенного текста.",
            "Для каждой истории той же embedding-моделью кодируются positive anchors и negative anchors. Positive anchors описывают широкий экономический контекст: inflation and cost of living, household finance and credit, labour and business activity, housing and mortgages.",
            "Negative anchors описывают семантически близкий шум: sport, crime and emergency reporting, arts and reviews, advertorial and property listing, generic local community stories.",
        ],
        Cm(1.2),
        Cm(4.3),
        Cm(16.0),
        Cm(11.0),
        font_size=13,
    )
    add_paragraph_box(
        slide,
        [
            "Для истории i считаются три величины:",
            "econ_positive_max(i) = max cosine(story_i, positive_anchor_j)",
            "econ_negative_max(i) = max cosine(story_i, negative_anchor_k)",
            "econ_margin(i) = econ_positive_max(i) - econ_negative_max(i)",
            "История считается экономической, если positive similarity достаточно высока, положительный смысл выигрывает у отрицательного, а negative anchor не доминирует.",
            "Strong keyword помогает только через мягкое ослабление порога и не даёт automatic pass.",
        ],
        Cm(18.0),
        Cm(4.9),
        Cm(13.8),
        Cm(10.4),
        font_size=13,
        fill=True,
        title="Критерий решения is_econ",
        title_color=ACCENT_BLUE,
    )


def build_slide_7_topics(prs: Presentation) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title(slide, "Как присваивались темы Bank of England", "Темы задаются как явные семантические профили, а не как одиночные ключевые слова.")
    add_paragraph_box(
        slide,
        [
            "Каждая тема BoE описывается в topic_profiles.yaml через четыре компоненты: name, description, keywords и seed_examples.",
            "Из этих частей собирается единый текст темы. Затем этот текст кодируется в topic embedding той же моделью, что и истории.",
            "Для каждой экономической истории считаются cosine similarity ко всем 7 topic embeddings. Далее выбираются top_topic, top1_similarity и top2_similarity.",
            "Если разница между top1_similarity и top2_similarity слишком мала, история помечается как topic_ambiguous.",
        ],
        Cm(1.2),
        Cm(4.6),
        Cm(16.0),
        Cm(10.2),
        font_size=13,
    )
    add_paragraph_box(
        slide,
        [
            "Присвоение темы идёт в два шага.",
            "Если история не прошла economic gate, assigned_topic = unassigned.",
            "Если история экономическая и top1_similarity превышает topic threshold, assigned_topic = top_topic.",
            "Если история экономическая, но ни одна из 7 тем не набирает достаточную similarity, assigned_topic = other_econ.",
            "Это решение сохраняет широкий экономический фон и не подменяет economic relevance близостью к заранее заданному набору тем.",
        ],
        Cm(18.0),
        Cm(4.9),
        Cm(13.8),
        Cm(10.0),
        font_size=13,
        fill=True,
        title="Итоговое назначение темы",
        title_color=ACCENT_BLUE,
    )


def build_slide_8_risk(prs: Presentation) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title(slide, "Как считался региональный индекс риска", "Индекс измеряет не только объём публикаций, но и отклонение темы от собственной локальной истории.")
    add_paragraph_box(
        slide,
        [
            "После story-level scoring результаты переносятся обратно на article-level строки. Для каждой статьи сохраняются is_econ, assigned_topic, top1_similarity, dup_weight и region_spread.",
            "Далее считается effective_weight = dup_weight / sqrt(region_spread). Эта величина понижает вклад широко синдицированных историй, которые одновременно опубликованы во многих регионах.",
            "Для каждой ячейки region x week x topic считаются y — сумма effective_weight по теме, и N — total_econ_weight, то есть сумма effective_weight по всем экономическим статьям региона на этой неделе.",
            "Сырая доля темы равна topic_share = y / N. Затем вводится сглаженная оценка p_post = (y + alpha_t) / (N + alpha_t + beta_t), где alpha_t и beta_t задаются через глобальную долю темы и prior_strength.",
        ],
        Cm(1.2),
        Cm(4.3),
        Cm(15.8),
        Cm(11.0),
        font_size=12,
    )
    add_paragraph_box(
        slide,
        [
            "baseline_share — среднее p_post за предыдущие 8 недель. recent_share — среднее p_post за предыдущие 4 недели.",
            "surprise = p_post - baseline_share показывает, насколько тема отклонилась от обычного уровня региона.",
            "momentum = p_post - recent_share показывает краткосрочное ускорение темы.",
            "mean_similarity — взвешенная средняя semantic similarity по статьям темы.",
            "support = sqrt(N / (N + 5)) понижает доверие к ячейкам со слабым покрытием.",
            "Итог: risk_score = support * (0.5 * z(surprise) + 0.3 * z(momentum) + 0.2 * z(mean_similarity)).",
        ],
        Cm(18.0),
        Cm(4.6),
        Cm(13.8),
        Cm(10.8),
        font_size=12,
        fill=True,
        title="Компоненты индекса",
        title_color=ACCENT_BLUE,
    )


def build_slide_9_results(prs: Presentation, hot_pairs: list[dict[str, str]]) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title(slide, "Результаты I: сильнейшие сочетания регионов и тем", "Публичная версия сохраняет только агрегированные регионально-тематические результаты без текстовых цитат из корпуса.")
    rows = [["Неделя", "Регион", "Тема", "Risk score"]]
    for row in hot_pairs[:6]:
        rows.append([
            row.get("week", ""),
            row.get("region", ""),
            row.get("topic", ""),
            f"{float(row.get('risk_score', 0.0)):.2f}",
        ])
    if len(rows) == 1:
        rows.append(["n/a", "n/a", "n/a", "n/a"])
    add_table(slide, rows, Cm(1.2), Cm(4.8), Cm(19.4), Cm(8.3))
    add_paragraph_box(
        slide,
        [
            "Внутри GitHub-версии репозитория не публикуются representative snippets и review-таблицы с прямыми фрагментами корпуса.",
            "Для интерпретации в публичной версии используются агрегированные пары region-topic, итоговые карты и тренды.",
            "Полные локальные диагностические таблицы остаются внешними артефактами и не входят в публичный контур репозитория.",
        ],
        Cm(21.1),
        Cm(4.9),
        Cm(11.2),
        Cm(7.8),
        font_size=13,
        fill=True,
        title="Публичный режим интерпретации",
        title_color=ACCENT_BLUE,
    )


def build_slide_10_maps(prs: Presentation) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title(slide, "Результаты II: карты", "Карты показывают как общий уровень регионального риска, так и доминирующую тему в каждом регионе.")
    add_picture_fit(slide, FIGURES_DIR / "overall_top_risk_map.png", Cm(1.2), Cm(4.6), Cm(15.2), Cm(11.6))
    add_picture_fit(slide, FIGURES_DIR / "dominant_topic_map.png", Cm(17.3), Cm(4.6), Cm(15.2), Cm(11.6))
    add_paragraph_box(
        slide,
        [
            "Слева: общий top-risk уровень региона.",
            "Справа: тема BoE, которая даёт максимальный риск в регионе.",
        ],
        Cm(1.2),
        Cm(16.4),
        Cm(31.3),
        Cm(1.5),
        font_size=12,
    )


def build_slide_11_trends(prs: Presentation) -> None:
    slide = prs.slides.add_slide(blank_layout(prs))
    add_title(slide, "Результаты III: динамика, стек и ограничения", "Тренды показывают, как меняется совокупный региональный риск во времени.")
    add_picture_fit(slide, FIGURES_DIR / "overall_region_risk_trends.png", Cm(1.2), Cm(4.8), Cm(19.0), Cm(10.0))
    add_picture_fit(slide, FIGURES_DIR / "hot_topic_trends.png", Cm(20.8), Cm(4.8), Cm(11.8), Cm(6.0))
    add_paragraph_box(
        slide,
        [
            "Стек: Python, pandas, numpy, pyarrow, sentence-transformers, matplotlib, seaborn, python-pptx, YAML.",
            "Ограничения: текущий прогон является sampled semantic pilot; ручная gold-label калибровка пока не выполнена; результат чувствителен к пороговым настройкам и параметрам сглаживания.",
            "Следующий шаг: полный локальный прогон по корпусу, ручная validation sample и сопоставление региональных сигналов с реальными коммуникационными событиями.",
        ],
        Cm(20.8),
        Cm(11.2),
        Cm(11.8),
        Cm(5.4),
        font_size=11,
        fill=True,
        title="Интерпретация",
        title_color=ACCENT_BLUE,
    )


def build_presentation() -> Path:
    prs = load_presentation()
    remove_all_slides(prs)

    summary = read_cleaning_summary()
    hot_pairs = read_hot_pairs()

    build_slide_1_title(prs)
    build_slide_2_problem(prs)
    build_slide_3_dataset_source(prs)
    build_slide_4_dataset_coverage(prs, summary)
    build_slide_5_cleaning(prs)
    build_slide_6_econ_gate(prs)
    build_slide_7_topics(prs)
    build_slide_8_risk(prs)
    build_slide_9_results(prs, hot_pairs)
    build_slide_10_maps(prs)
    build_slide_11_trends(prs)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(OUTPUT_PATH)
        return OUTPUT_PATH
    except PermissionError:
        prs.save(FALLBACK_OUTPUT_PATH)
        return FALLBACK_OUTPUT_PATH


def main() -> None:
    output_path = build_presentation()
    print(f"Presentation saved to: {output_path}")


if __name__ == "__main__":
    main()
